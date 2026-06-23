"""Génère la présentation PowerPoint TokenVeil (RSSI).

N'utilise que des APIs officielles python-pptx (pas de manipulation XML
manuelle) : une version précédente bricolait l'ordre des éléments OOXML pour
le fond de slide et les flèches, ce qui corrompait le fichier (PowerPoint
demandait une réparation à l'ouverture).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

ACCENT = RGBColor(0xD9, 0x77, 0x57)
ACCENT_DARK = RGBColor(0xB8, 0x5A, 0x3D)
DARK = RGBColor(0x2B, 0x29, 0x24)
DARK_SOFT = RGBColor(0x3D, 0x2F, 0x27)
LIGHT_BG = RGBColor(0xFA, 0xF8, 0xF5)
TEXT_DIM = RGBColor(0x6B, 0x64, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF1, 0xEC, 0xE3)
LINE_SOFT = RGBColor(0xE5, 0xDE, 0xD3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg=LIGHT_BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def add_text(slide, left, top, width, height, text, size=18, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, anchor=None, italic=False,
             line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK,
                 space_after=10, sub_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    sub_color = sub_color or TEXT_DIM
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        sub = isinstance(item, tuple)
        run = p.add_run()
        if sub:
            run.text = f"   –  {item[0]}"
            run.font.size = Pt(size - 2)
            run.font.color.rgb = sub_color
        else:
            run.text = f"●  {item}"
            run.font.size = Pt(size)
            run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def kicker(slide, text, dark_bg=False):
    add_text(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.4), text,
              size=13, color=ACCENT, bold=True)


def title(slide, text, color=DARK, size=30):
    add_text(slide, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.8), text,
              size=size, color=color, bold=True)


def footer(slide, n, dark_bg=False):
    color = RGBColor(0x9B, 0x92, 0x86) if dark_bg else TEXT_DIM
    add_text(slide, Inches(0.7), Inches(7.08), Inches(7), Inches(0.32),
              "TokenVeil — Alpha interne", size=10, color=color)
    add_text(slide, Inches(12.0), Inches(7.08), Inches(0.6), Inches(0.32),
              str(n), size=10, color=color, align=PP_ALIGN.RIGHT)


def node(slide, left, top, w, h, text, fill=WHITE, text_color=DARK, size=12,
         bold=True, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1.25)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = "Calibri"
    return sh


def connector(slide, x1, y1, x2, y2, color=TEXT_DIM, width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def flow_arrow_label(slide, cx, cy, color=ACCENT):
    # flèche simple en texte (évite toute manipulation XML des terminaisons
    # de trait, source de corruption du fichier dans une version antérieure)
    add_text(slide, cx, cy, Inches(0.3), Inches(0.3), "›", size=20,
              color=color, bold=True, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)


# ============================================================ SLIDE 1 — TITRE
s = add_slide(bg=DARK)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.55), Inches(0.18), Inches(1.4))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
bar.shadow.inherit = False
add_text(s, Inches(0.75), Inches(2.7), Inches(11), Inches(1.0), "TokenVeil",
          size=46, color=WHITE, bold=True)
add_text(s, Inches(0.78), Inches(3.55), Inches(11), Inches(0.7),
          "Discuter avec Claude sans jamais lui exposer une donnée sensible",
          size=20, color=RGBColor(0xD8, 0xD2, 0xC6))
add_text(s, Inches(0.78), Inches(4.3), Inches(11), Inches(0.5),
          "Présentation RSSI  ·  Build Alpha  ·  19 juin 2026", size=14, color=ACCENT, bold=True)
add_text(s, Inches(0.78), Inches(6.7), Inches(6), Inches(0.4),
          "Marc Bourrel — DevOps & Cloud Architect", size=12, color=TEXT_DIM)

# ============================================================ SLIDE 2 — CONSTAT
s = add_slide()
kicker(s, "LE CONSTAT")
title(s, "Vos équipes collent déjà des données sensibles dans une IA")
risks = [
    "Pour diagnostiquer un incident, on copie un log technique dans Claude ou ChatGPT — et avec lui, des IP internes, des clés API, parfois des identifiants de session",
    "Pour reformuler un mail ou analyser un cas client, on y colle aussi des données personnelles : noms, IBAN, références de contrat",
    "Ces données partent chez un prestataire tiers, sortent du périmètre de l'entreprise, et personne ne sait vraiment ce qu'il en advient",
    "On ne peut pas compter sur la seule sensibilisation pour empêcher ce réflexe — il faut une barrière technique",
]
add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.2), risks, size=18, space_after=22)
footer(s, 2)

# ============================================================ SLIDE 3 — PRINCIPE
s = add_slide()
kicker(s, "LE PRINCIPE")
title(s, "On anonymise avant l'envoi, on restaure à l'affichage")
add_text(s, Inches(0.8), Inches(1.78), Inches(11.5), Inches(0.5),
          "TokenVeil s'intercale entre l'utilisateur et Claude — sans rien changer à son usage.",
          size=15, color=TEXT_DIM)

steps = [
    ("1", "L'utilisateur colle un vrai\nlog ou une vraie donnée"),
    ("2", "Le moteur remplace tout\npar des tokens opaques\n<IP_ADDRESS_1>, <API_SECRET_2>..."),
    ("3", "Claude ne reçoit\nque ces tokens"),
    ("4", "La réponse est restaurée\nau moment de l'affichage"),
]
x0 = Inches(0.8); w = Inches(2.75); gap = Inches(0.3); y = Inches(2.75); h = Inches(1.85)
for i, (n, txt) in enumerate(steps):
    bx = x0 + i * (w + gap)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, bx + Inches(1.1), y - Inches(0.5), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background(); circ.shadow.inherit = False
    ctf = circ.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = n; cr.font.size = Pt(16); cr.font.bold = True; cr.font.color.rgb = WHITE
    node(s, bx, y, w, h, txt, fill=WHITE, text_color=DARK, size=13, bold=False, line_color=LINE_SOFT)
    if i < len(steps) - 1:
        flow_arrow_label(s, bx + w + Inches(0.02), y + h / 2 - Inches(0.15), color=ACCENT)

add_text(s, Inches(0.8), Inches(5.05), Inches(11.5), Inches(0.7),
          "La donnée réelle ne quitte jamais le serveur interne : Anthropic ne voit et ne renvoie que des identifiants opaques.",
          size=15, color=ACCENT, bold=True)
footer(s, 3)

# ============================================================ SLIDE 4 — ARCHITECTURE
s = add_slide(bg=DARK)
kicker(s, "ARCHITECTURE", dark_bg=True)
title(s, "Le flux de données, étape par étape", color=WHITE)

row_y = Inches(2.75)
row_h = Inches(1.15)
boxes = [
    ("Utilisateur\n(navigateur)", Inches(0.5), Inches(2.0), DARK_SOFT, WHITE),
    ("Interface web\n(FastAPI + auth LDAP)", Inches(2.7), Inches(2.3), DARK_SOFT, WHITE),
    ("Moteur\nd'anonymisation", Inches(5.2), Inches(2.1), ACCENT, WHITE),
    ("Claude Code CLI\n(OAuth, abonnement perso.)", Inches(7.5), Inches(2.5), DARK_SOFT, WHITE),
    ("api.anthropic.com\n(ne voit que des tokens)", Inches(10.2), Inches(2.6), RGBColor(0x4A, 0x44, 0x3C), WHITE),
]
positions = []
for label, left, w, fill, tc in boxes:
    node(s, left, row_y, w, row_h, label, fill=fill, text_color=tc, size=12.5)
    positions.append((left, w))

cursor = positions[0][0] + positions[0][1]
for i in range(1, len(positions)):
    flow_arrow_label(s, cursor + Inches(0.02), row_y + row_h / 2 - Inches(0.15))
    cursor = positions[i][0] + positions[i][1]

# boîte "stockage" au-dessus du moteur d'anonymisation, avec une marge nette
# sous le titre pour ne jamais chevaucher (cause de la corruption précédente)
db_top = Inches(1.55)
db_left, db_w = Inches(5.05), Inches(2.4)
node(s, db_left, db_top, db_w, Inches(0.85),
     "SQLite chiffrée — messages\nanonymisés, mapping Fernet",
     fill=DARK_SOFT, text_color=WHITE, size=11.5)
mid_x = positions[2][0] + positions[2][1] / 2
connector(s, mid_x, db_top + Inches(0.85), mid_x, row_y, color=TEXT_DIM, width=1.5)

add_text(s, Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.45),
          "À partir du moteur d'anonymisation, plus aucune donnée réelle ne circule — uniquement des tokens.",
          size=13.5, color=RGBColor(0xD8, 0xD2, 0xC6), italic=True)

bottom_pts = [
    "Chaque utilisateur lie son propre abonnement Claude Pro/Max — pas de clé API partagée, pas de facturation au volume",
    "Le mapping entre un token et la vraie valeur reste chiffré au repos (Fernet) ; seul le propriétaire authentifié peut le déchiffrer",
    "L'authentification s'appuie sur le LDAP/Active Directory déjà en place, avec restriction possible par groupe",
]
add_bullets(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.3), bottom_pts, size=14.5, color=WHITE, space_after=14)
footer(s, 4, dark_bg=True)

# ============================================================ SLIDE 5 — MOTEUR (PRESIDIO)
s = add_slide()
kicker(s, "LE MOTEUR D'ANONYMISATION")
title(s, "Sous le capot : Microsoft Presidio, spaCy, et nos propres règles")

in_box = node(s, Inches(4.5), Inches(1.55), Inches(4.3), Inches(0.6),
              "Ligne de log brute",
              fill=DARK_SOFT, text_color=WHITE, size=12.5)

left_box = node(s, Inches(0.9), Inches(2.45), Inches(5.0), Inches(1.05),
                "spaCy NER (fr_core_news_lg / en_core_web_lg)\nPersonnes, organisations, lieux...",
                fill=WHITE, text_color=DARK, size=11.5, line_color=LINE_SOFT)
right_box = node(s, Inches(7.4), Inches(2.45), Inches(5.0), Inches(1.05),
                  "Recognizers regex maison\nIP, clés API, IBAN, hostnames, mots-clés métier",
                  fill=WHITE, text_color=DARK, size=11.5, line_color=LINE_SOFT)
connector(s, Inches(6.65), Inches(2.15), Inches(3.4), Inches(2.45), color=LINE_SOFT, width=1.5)
connector(s, Inches(6.65), Inches(2.15), Inches(9.9), Inches(2.45), color=LINE_SOFT, width=1.5)

merge_box = node(s, Inches(3.4), Inches(3.75), Inches(6.2), Inches(0.65),
                  "Résolution des chevauchements — le score de confiance le plus élevé gagne",
                  fill=DARK_SOFT, text_color=WHITE, size=11.5)
connector(s, Inches(3.4), Inches(3.5), Inches(5.5), Inches(3.75), color=LINE_SOFT, width=1.5)
connector(s, Inches(9.9), Inches(3.5), Inches(7.3), Inches(3.75), color=LINE_SOFT, width=1.5)

final_box = node(s, Inches(3.9), Inches(4.65), Inches(5.2), Inches(0.6),
                  "Texte tokenisé envoyé à Claude",
                  fill=ACCENT, text_color=WHITE, size=12.5)
connector(s, Inches(6.5), Inches(4.4), Inches(6.5), Inches(4.65), color=LINE_SOFT, width=1.5)

bottom_pts5 = [
    "Presidio (Microsoft, open source) fournit le moteur d'analyse ; spaCy fait la reconnaissance d'entités généraliste",
    "Nos propres règles couvrent ce que la détection générique ne peut pas savoir : secrets, hostnames, mots-clés propres à l'entreprise",
]
add_bullets(s, Inches(0.8), Inches(5.55), Inches(11.5), Inches(1.4), bottom_pts5, size=13.5, space_after=10)
footer(s, 5)

# ============================================================ SLIDE 6 — DETECTION
s = add_slide()
kicker(s, "CE QUI EST DÉTECTÉ")
title(s, "Une couverture pensée pour des vrais logs techniques")
rows = [
    ("Réseau", "Adresses IPv4, adresses MAC, hostnames internes (.local / .corp / .lan)"),
    ("Identité", "Noms, emails, téléphones, champs de logs user=/login=/owner="),
    ("Secrets", "Clés API, tokens, mots de passe, bearer tokens, chaînes haute-entropie"),
    ("Financier", "IBAN, numéros de carte bancaire"),
    ("Références métier", "Références client, ticket ou employé (CUST-1234, TICKET-5678...)"),
    ("Détection générique", "Organisations, lieux, numéros de sécurité sociale (FR et EN)"),
]
top = Inches(1.95)
rh = Inches(0.75)
for i, (cat, ex) in enumerate(rows):
    y = top + i * rh
    fill = WHITE if i % 2 == 0 else ROW_ALT
    rowbg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.5), rh)
    rowbg.fill.solid(); rowbg.fill.fore_color.rgb = fill; rowbg.line.fill.background(); rowbg.shadow.inherit = False
    add_text(s, Inches(0.95), y + Inches(0.13), Inches(2.6), Inches(0.5), cat, size=15, bold=True, color=ACCENT_DARK)
    add_text(s, Inches(3.7), y + Inches(0.13), Inches(8.4), Inches(0.55), ex, size=14, color=DARK)
footer(s, 6)

# ============================================================ SLIDE 7 — MOTS-CLES & AUTO-DETECTION
s = add_slide()
kicker(s, "PERSONNALISATION")
title(s, "Les mots que seule l'entreprise connaît")
left_items7 = [
    "Un nom de projet interne, un codename, le petit nom d'un service maison : la détection générique ne peut pas les connaître à l'avance",
    "Ajoutés une fois dans Réglages, ils sont systématiquement remplacés par un token — pour tous les utilisateurs, sur tous les messages",
    "Toujours prioritaires : un mot-clé métier l'emporte sur n'importe quelle détection générique concurrente",
]
add_bullets(s, Inches(0.8), Inches(1.95), Inches(5.5), Inches(4.6), left_items7, size=15, space_after=20)

panel7 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.95), Inches(5.6), Inches(4.5))
panel7.fill.solid(); panel7.fill.fore_color.rgb = DARK; panel7.line.fill.background(); panel7.shadow.inherit = False
add_text(s, Inches(7.1), Inches(2.2), Inches(5.0), Inches(0.5),
          "Auto-détection, sans rien exposer", size=15, bold=True, color=ACCENT)
verif7 = [
    "Un script autonome (aucune dépendance) à lancer sur le serveur à protéger — jamais sur celui qui héberge la webapp",
    "Il lit uniquement les noms de containers Docker et /etc/hosts ; jamais un fichier de config ou un secret",
    "Le résultat se colle dans l'interface : rien n'est ajouté avant qu'un humain n'ait coché et validé chaque terme",
    "Zéro nouveau privilège côté webapp : pas de socket Docker monté, pas d'accès réseau supplémentaire",
]
add_bullets(s, Inches(7.1), Inches(2.75), Inches(5.0), Inches(3.5), verif7, size=13.5, color=WHITE, space_after=15)
footer(s, 7)

# ============================================================ SLIDE 8 — SECURITE
s = add_slide()
kicker(s, "SÉCURITÉ DES DONNÉES")
title(s, "Ce qu'on stocke, et comment c'est protégé")
left_items = [
    "La base de conversations ne contient que la version anonymisée des messages — le texte réel n'est jamais écrit sur disque",
    "Le mapping token ↔ valeur réelle est chiffré au repos (Fernet), et n'est déchiffré en mémoire que pour le propriétaire authentifié, à la demande",
    "Les jetons OAuth Claude sont chiffrés par utilisateur, avec des droits fichier restreints, et ne sont jamais journalisés",
]
add_bullets(s, Inches(0.8), Inches(1.95), Inches(5.5), Inches(4.6), left_items, size=15, space_after=20)

panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.95), Inches(5.6), Inches(4.5))
panel.fill.solid(); panel.fill.fore_color.rgb = DARK; panel.line.fill.background(); panel.shadow.inherit = False
add_text(s, Inches(7.1), Inches(2.2), Inches(5.0), Inches(0.5),
          "Testé sur des données réelles", size=15, bold=True, color=ACCENT)
verif = [
    "Tests menés sur de vrais logs d'infrastructure (Sonarr, Nextcloud, Jellyfin, syslog systemd)",
    "Une clé API qui fuitait en clair dans un log a bien été tokenisée",
    "Les IP internes et publiques ont bien été tokenisées",
    "Vérification directe dans la base SQLite brute : aucune fuite en clair constatée",
]
add_bullets(s, Inches(7.1), Inches(2.75), Inches(5.0), Inches(3.5), verif, size=13.5, color=WHITE, space_after=15)
footer(s, 8)

# ============================================================ SLIDE 7 — GOUVERNANCE
s = add_slide()
kicker(s, "GOUVERNANCE & ACCÈS")
title(s, "Une authentification d'entreprise, une responsabilité individuelle")
items = [
    "Connexion via le LDAP ou l'Active Directory déjà en place — pas de nouvel annuaire à gérer",
    ("Restriction possible à un groupe LDAP dédié", 1),
    "Chaque utilisateur lie son propre abonnement Claude Pro/Max depuis l'interface, en quelques clics",
    ("Aucune clé API partagée et facturée au volume : l'usage de chacun reste rattaché à son abonnement personnel", 1),
    "On sait qui a posé quelle question — sans que la donnée réelle n'ait jamais été exposée à Claude",
]
add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), items, size=17, space_after=20)
footer(s, 9)

# ============================================================ SLIDE 8 — TRANSPARENCE
s = add_slide()
kicker(s, "TRANSPARENCE UTILISATEUR")
title(s, "L'utilisateur voit lui-même ce qui part réellement")
items = [
    "Un aperçu s'affiche en temps réel pendant la frappe : la version anonymisée du message apparaît avant même l'envoi",
    "Sur chaque message envoyé, un bouton « Voir la version envoyée à Claude » révèle le texte tokenisé exact qui a quitté le serveur",
    "Rien n'est caché : la promesse d'anonymisation se vérifie à l'œil, message après message",
]
add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5), items, size=17, space_after=22)
footer(s, 10)

# ============================================================ SLIDE 9 — STATUT / ROADMAP
s = add_slide()
kicker(s, "STATUT")
title(s, "Une alpha qui tient déjà la route, avec une suite claire")
add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
          "Le mécanisme central est validé de bout en bout. Voici ce qu'il reste à construire avant un déploiement plus large :",
          size=15, color=TEXT_DIM)
roadmap = [
    "Isolation stricte des conversations entre utilisateurs",
    "Journal d'audit des décisions d'anonymisation, pour les revues de conformité",
    "Politiques d'accès multi-tenant basées sur les groupes LDAP",
]
add_bullets(s, Inches(0.8), Inches(2.45), Inches(11.5), Inches(3.5), roadmap, size=16, space_after=18)
footer(s, 11)

# ============================================================ SLIDE 10 — DEMO
s = add_slide(bg=DARK)
add_text(s, Inches(0.78), Inches(3.0), Inches(11.5), Inches(1.0), "Place à la démonstration",
          size=38, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(3.9), Inches(11), Inches(0.6),
          "Quelques vrais-faux logs envoyés en direct, pour voir l'anonymisation se faire en temps réel",
          size=18, color=ACCENT)
footer(s, 12, dark_bg=True)

prs.save("/home/kemar/docker/tokenveil/presentation/TokenVeil_RSSI.pptx")
print("OK")
